"""
Background worker for processing document ingestion.
This module handles document parsing, chunking, embedding generation,
and knowledge graph creation for uploaded files.
"""

import asyncio
import logging
import json
import hashlib
import mimetypes
from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime
import traceback
import io
import re

# Document processing libraries
import pypdf
from docx import Document
import pandas as pd
import openpyxl
from pptx import Presentation
from PIL import Image
import easyocr
from sentence_transformers import SentenceTransformer

# NLP libraries
import spacy
from spacy import displacy
import numpy as np

from ..core.db_mongo import get_mongodb_client, MongoDBClient
from ..core.db_qdrant import get_qdrant_client, QdrantDBClient
from ..core.db_arango import get_arango_client, ArangoDBClient
from ..core.cache_redis import get_redis_client, RedisClient
from ..config import settings

logger = logging.getLogger(__name__)

class IngestionWorker:
    """Worker for processing uploaded documents in background."""
    
    def __init__(self):
        self.embedding_model = None
        self.nlp_model = None
        self.ocr_reader = None
        self.mongo_client = None
        self.qdrant_client = None
        self.arango_client = None
        self.redis_client = None
        self.is_processing = False
    
    async def initialize(self):
        """Initialize the ingestion worker."""
        try:
            logger.info("Initializing ingestion worker...")
            
            # Get database clients
            self.mongo_client = await get_mongodb_client()
            self.qdrant_client = await get_qdrant_client()
            self.arango_client = await get_arango_client()
            self.redis_client = await get_redis_client()
            
            # Load embedding model
            logger.info(f"Loading embedding model: {settings.embedding_model_name}")
            self.embedding_model = SentenceTransformer(settings.embedding_model_name)
            
            # Load NLP model for entity extraction
            try:
                logger.info("Loading spaCy model for entity extraction...")
                self.nlp_model = spacy.load("en_core_web_sm")
            except IOError:
                logger.warning("spaCy model not found. Installing...")
                # In production, this should be handled during setup
                self.nlp_model = None
            
            # Initialize OCR reader for image processing
            if settings.enable_ocr:
                try:
                    logger.info("Initializing OCR reader...")
                    self.ocr_reader = easyocr.Reader(['en'])
                except Exception as e:
                    logger.warning(f"OCR initialization failed: {e}")
                    self.ocr_reader = None
            
            logger.info("Ingestion worker initialized successfully")
            
        except Exception as e:
            logger.error(f"Failed to initialize ingestion worker: {e}")
            raise
    
    async def process_document(self, doc_id: str) -> Dict[str, Any]:
        """
        Process a single document through the ingestion pipeline.
        
        Args:
            doc_id: Document ID to process
        
        Returns:
            Processing result with status and stats
        """
        start_time = datetime.utcnow()
        
        try:
            # Mark as processing
            self.is_processing = True
            
            # Update status
            await self._update_processing_status(doc_id, "processing", "Starting document processing...")
            
            # Get document from MongoDB
            document = await self.mongo_client.get_document(doc_id)
            if not document:
                raise ValueError(f"Document {doc_id} not found")
            
            logger.info(f"Processing document: {document['filename']} ({document['file_type']})")
            
            # Step 1: Extract text content
            await self._update_processing_status(doc_id, "processing", "Extracting text content...")
            text_content = await self._extract_text_content(document)
            
            if not text_content:
                raise ValueError("No text content could be extracted from document")
            
            # Step 2: Clean and preprocess text
            await self._update_processing_status(doc_id, "processing", "Preprocessing text...")
            cleaned_text = self._clean_text(text_content)
            
            # Step 3: Create text chunks
            await self._update_processing_status(doc_id, "processing", "Creating text chunks...")
            chunks = self._create_chunks(cleaned_text, document['filename'])
            
            logger.info(f"Created {len(chunks)} chunks for document {doc_id}")
            
            # Step 4: Generate embeddings
            await self._update_processing_status(doc_id, "processing", "Generating embeddings...")
            embeddings = await self._generate_embeddings(chunks)
            
            # Step 5: Store chunks in vector database
            await self._update_processing_status(doc_id, "processing", "Storing vectors...")
            vector_ids = await self._store_vectors(doc_id, chunks, embeddings)
            
            # Step 6: Extract entities and relationships for knowledge graph
            entities_extracted = 0
            relationships_extracted = 0
            
            if settings.enable_knowledge_graph and self.nlp_model:
                await self._update_processing_status(doc_id, "processing", "Extracting entities...")
                entities, relationships = await self._extract_entities_and_relationships(cleaned_text, doc_id)
                
                if entities:
                    await self.arango_client.upsert_entities(entities)
                    entities_extracted = len(entities)
                
                if relationships:
                    await self.arango_client.create_relationships(relationships)
                    relationships_extracted = len(relationships)
            
            # Step 7: Update document metadata
            processing_stats = {
                "chunks_created": len(chunks),
                "vectors_stored": len(vector_ids),
                "entities_extracted": entities_extracted,
                "relationships_extracted": relationships_extracted,
                "processed_at": datetime.utcnow(),
                "processing_time_seconds": (datetime.utcnow() - start_time).total_seconds(),
                "text_length": len(cleaned_text),
                "vector_ids": vector_ids
            }
            
            await self.mongo_client.update_document_processing_status(
                doc_id, "completed", processing_stats
            )
            
            # Update status cache
            await self._update_processing_status(doc_id, "completed", "Document processed successfully")
            
            logger.info(f"Document {doc_id} processed successfully in {processing_stats['processing_time_seconds']:.2f}s")
            
            # Notify via Redis pub/sub
            await self.redis_client.publish_processing_update({
                "doc_id": doc_id,
                "status": "completed",
                "stats": processing_stats
            })
            
            return {
                "status": "success",
                "doc_id": doc_id,
                "stats": processing_stats
            }
            
        except Exception as e:
            error_msg = f"Processing failed: {str(e)}"
            logger.error(f"Document {doc_id} processing failed: {e}")
            logger.error(traceback.format_exc())
            
            # Update error status
            await self.mongo_client.update_document_processing_status(
                doc_id, "failed", {"error": error_msg, "failed_at": datetime.utcnow()}
            )
            
            await self._update_processing_status(doc_id, "failed", error_msg)
            
            # Notify error
            await self.redis_client.publish_processing_update({
                "doc_id": doc_id,
                "status": "failed",
                "error": error_msg
            })
            
            return {
                "status": "error",
                "doc_id": doc_id,
                "error": error_msg
            }
            
        finally:
            self.is_processing = False
    
    async def _extract_text_content(self, document: Dict[str, Any]) -> str:
        """Extract text content based on file type."""
        file_type = document['file_type']
        filename = document['filename']
        
        # Get file content from GridFS
        file_content = await self.mongo_client.get_file_content(document['_id'])
        if not file_content:
            raise ValueError("Could not retrieve file content")
        
        try:
            if file_type == 'pdf':
                return await self._extract_pdf_text(file_content)
            elif file_type in ['doc', 'docx']:
                return await self._extract_docx_text(file_content)
            elif file_type in ['txt', 'md']:
                return file_content.decode('utf-8')
            elif file_type == 'csv':
                return await self._extract_csv_text(file_content)
            elif file_type in ['xls', 'xlsx']:
                return await self._extract_excel_text(file_content)
            elif file_type in ['ppt', 'pptx']:
                return await self._extract_pptx_text(file_content)
            elif file_type in ['jpg', 'jpeg', 'png', 'tiff', 'bmp']:
                return await self._extract_image_text(file_content)
            elif file_type == 'json':
                return await self._extract_json_text(file_content)
            else:
                # Try to decode as text
                try:
                    return file_content.decode('utf-8')
                except UnicodeDecodeError:
                    raise ValueError(f"Unsupported file type: {file_type}")
                    
        except Exception as e:
            logger.error(f"Text extraction failed for {filename}: {e}")
            raise
    
    async def _extract_pdf_text(self, content: bytes) -> str:
        """Extract text from PDF content."""
        try:
            pdf_file = io.BytesIO(content)
            reader = pypdf.PdfReader(pdf_file)
            
            text_parts = []
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                except Exception as e:
                    logger.warning(f"Could not extract text from page {page_num + 1}: {e}")
                    continue
            
            return "\n\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise ValueError(f"Could not extract text from PDF: {e}")
    
    async def _extract_docx_text(self, content: bytes) -> str:
        """Extract text from DOCX content."""
        try:
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise ValueError(f"Could not extract text from DOCX: {e}")
    
    async def _extract_csv_text(self, content: bytes) -> str:
        """Extract text from CSV content."""
        try:
            csv_file = io.StringIO(content.decode('utf-8'))
            df = pd.read_csv(csv_file)
            
            # Convert DataFrame to readable text
            text_parts = [f"CSV Data with {len(df)} rows and {len(df.columns)} columns"]
            text_parts.append(f"Columns: {', '.join(df.columns)}")
            
            # Add first few rows as examples
            for idx, row in df.head(10).iterrows():
                row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                text_parts.append(f"Row {idx + 1}: {row_text}")
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"CSV extraction failed: {e}")
            raise ValueError(f"Could not extract text from CSV: {e}")
    
    async def _extract_excel_text(self, content: bytes) -> str:
        """Extract text from Excel content."""
        try:
            excel_file = io.BytesIO(content)
            workbook = openpyxl.load_workbook(excel_file)
            
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                
                # Get data from sheet
                for row in sheet.iter_rows(max_row=50, values_only=True):  # Limit to first 50 rows
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        text_parts.append(" | ".join(row_data))
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            raise ValueError(f"Could not extract text from Excel: {e}")
    
    async def _extract_pptx_text(self, content: bytes) -> str:
        """Extract text from PowerPoint content."""
        try:
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = [f"[Slide {slide_num + 1}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise ValueError(f"Could not extract text from PPTX: {e}")
    
    async def _extract_image_text(self, content: bytes) -> str:
        """Extract text from image using OCR."""
        if not self.ocr_reader:
            raise ValueError("OCR not available")
        
        try:
            image = Image.open(io.BytesIO(content))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Perform OCR
            results = self.ocr_reader.readtext(np.array(image))
            
            # Extract text with confidence scores
            text_parts = []
            for (bbox, text, confidence) in results:
                if confidence > 0.5:  # Only include text with reasonable confidence
                    text_parts.append(text)
            
            return "\n".join(text_parts) if text_parts else "No text detected in image"
            
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            raise ValueError(f"Could not extract text from image: {e}")
    
    async def _extract_json_text(self, content: bytes) -> str:
        """Extract readable text from JSON content."""
        try:
            data = json.loads(content.decode('utf-8'))
            
            def json_to_text(obj, prefix=""):
                """Convert JSON object to readable text."""
                if isinstance(obj, dict):
                    text_parts = []
                    for key, value in obj.items():
                        if isinstance(value, (dict, list)):
                            text_parts.append(f"{prefix}{key}:")
                            text_parts.append(json_to_text(value, prefix + "  "))
                        else:
                            text_parts.append(f"{prefix}{key}: {value}")
                    return "\n".join(text_parts)
                elif isinstance(obj, list):
                    text_parts = []
                    for i, item in enumerate(obj):
                        text_parts.append(f"{prefix}[{i}]: {json_to_text(item, prefix + '  ')}")
                    return "\n".join(text_parts)
                else:
                    return str(obj)
            
            return json_to_text(data)
            
        except Exception as e:
            logger.error(f"JSON extraction failed: {e}")
            raise ValueError(f"Could not extract text from JSON: {e}")
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text."""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s\.,;:!?()-]', '', text)
        
        # Fix common OCR errors
        text = text.replace('—', '-')
        text = text.replace('"', '"').replace('"', '"')
        text = text.replace(''', "'").replace(''', "'")
        
        # Remove excessive newlines
        text = re.sub(r'\n+', '\n', text)
        
        return text.strip()
    
    def _create_chunks(self, text: str, filename: str) -> List[Dict[str, Any]]:
        """Create overlapping text chunks for processing."""
        if not text:
            return []
        
        # Split into sentences first
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        chunks = []
        current_chunk = ""
        current_sentences = []
        
        for sentence in sentences:
            # Check if adding this sentence would exceed chunk size
            potential_chunk = current_chunk + " " + sentence if current_chunk else sentence
            
            if len(potential_chunk.split()) <= settings.chunk_size:
                current_chunk = potential_chunk
                current_sentences.append(sentence)
            else:
                # Save current chunk if it's not empty
                if current_chunk:
                    chunks.append({
                        "text": current_chunk.strip(),
                        "chunk_id": f"chunk_{len(chunks)}",
                        "sentence_count": len(current_sentences),
                        "word_count": len(current_chunk.split()),
                        "source_file": filename
                    })
                
                # Start new chunk with overlap
                if settings.chunk_overlap > 0 and len(current_sentences) > 1:
                    # Keep last few sentences for overlap
                    overlap_sentences = current_sentences[-settings.chunk_overlap:]
                    current_chunk = " ".join(overlap_sentences) + " " + sentence
                    current_sentences = overlap_sentences + [sentence]
                else:
                    current_chunk = sentence
                    current_sentences = [sentence]
        
        # Add the last chunk
        if current_chunk:
            chunks.append({
                "text": current_chunk.strip(),
                "chunk_id": f"chunk_{len(chunks)}",
                "sentence_count": len(current_sentences),
                "word_count": len(current_chunk.split()),
                "source_file": filename
            })
        
        return chunks
    
    async def _generate_embeddings(self, chunks: List[Dict[str, Any]]) -> List[List[float]]:
        """Generate embeddings for text chunks."""
        if not chunks:
            return []
        
        texts = [chunk["text"] for chunk in chunks]
        
        try:
            # Generate embeddings in batches to avoid memory issues
            batch_size = 32
            embeddings = []
            
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]
                batch_embeddings = self.embedding_model.encode(batch)
                embeddings.extend(batch_embeddings.tolist())
            
            return embeddings
            
        except Exception as e:
            logger.error(f"Embedding generation failed: {e}")
            raise
    
    async def _store_vectors(self, doc_id: str, chunks: List[Dict[str, Any]], 
                           embeddings: List[List[float]]) -> List[str]:
        """Store vectors in Qdrant database."""
        try:
            vector_data = []
            
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                vector_id = f"{doc_id}_{chunk['chunk_id']}"
                
                vector_data.append({
                    "id": vector_id,
                    "vector": embedding,
                    "payload": {
                        "doc_id": doc_id,
                        "chunk_id": chunk["chunk_id"],
                        "text": chunk["text"],
                        "source_file": chunk["source_file"],
                        "word_count": chunk["word_count"],
                        "sentence_count": chunk["sentence_count"],
                        "created_at": datetime.utcnow().isoformat()
                    }
                })
            
            # Store in batches
            await self.qdrant_client.upsert_vectors(vector_data)
            
            return [data["id"] for data in vector_data]
            
        except Exception as e:
            logger.error(f"Vector storage failed: {e}")
            raise
    
    async def _extract_entities_and_relationships(self, text: str, doc_id: str) -> tuple[List[Dict], List[Dict]]:
        """Extract entities and relationships for knowledge graph."""
        if not self.nlp_model:
            return [], []
        
        try:
            # Process text with spaCy
            doc = self.nlp_model(text[:1000000])  # Limit text length for processing
            
            entities = []
            relationships = []
            
            # Extract named entities
            for ent in doc.ents:
                if ent.label_ in ["PERSON", "ORG", "GPE", "PRODUCT", "EVENT", "MONEY", "DATE"]:
                    entity_data = {
                        "name": ent.text,
                        "type": ent.label_,
                        "doc_id": doc_id,
                        "start_char": ent.start_char,
                        "end_char": ent.end_char,
                        "confidence": 1.0  # spaCy doesn't provide confidence scores
                    }
                    entities.append(entity_data)
            
            # Extract relationships using dependency parsing
            for sent in doc.sents:
                for token in sent:
                    if token.dep_ in ["nsubj", "dobj", "pobj"]:
                        # Simple relationship extraction
                        if token.head.pos_ == "VERB":
                            subject = None
                            object_entity = None
                            
                            # Find subject and object
                            for child in token.head.children:
                                if child.dep_ == "nsubj":
                                    subject = child.text
                                elif child.dep_ in ["dobj", "pobj"]:
                                    object_entity = child.text
                            
                            if subject and object_entity:
                                relationship_data = {
                                    "from_entity": subject,
                                    "to_entity": object_entity,
                                    "relation_type": token.head.lemma_,
                                    "doc_id": doc_id,
                                    "confidence": 0.8
                                }
                                relationships.append(relationship_data)
            
            logger.info(f"Extracted {len(entities)} entities and {len(relationships)} relationships")
            return entities, relationships
            
        except Exception as e:
            logger.error(f"Entity extraction failed: {e}")
            return [], []
    
    async def _update_processing_status(self, doc_id: str, status: str, message: str):
        """Update processing status in Redis cache."""
        try:
            await self.redis_client.set(
                f"processing_status:{doc_id}",
                json.dumps({
                    "status": status,
                    "message": message,
                    "updated_at": datetime.utcnow().isoformat()
                }),
                ex=3600  # Expire after 1 hour
            )
        except Exception as e:
            logger.error(f"Failed to update processing status: {e}")
    
    async def get_processing_status(self, doc_id: str) -> Dict[str, Any]:
        """Get current processing status for a document."""
        try:
            status_data = await self.redis_client.get(f"processing_status:{doc_id}")
            if status_data:
                return json.loads(status_data)
            return {"status": "unknown", "message": "No status available"}
        except Exception as e:
            logger.error(f"Failed to get processing status: {e}")
            return {"status": "error", "message": "Could not retrieve status"}
    
    async def health_check(self) -> bool:
        """Health check for ingestion worker."""
        try:
            if not self.embedding_model:
                return False
            
            # Test embedding generation
            test_embedding = self.embedding_model.encode("test text")
            return len(test_embedding) == settings.embedding_dimension
            
        except Exception as e:
            logger.error(f"Ingestion worker health check failed: {e}")
            return False

# Global worker instance
ingestion_worker = IngestionWorker()

async def get_ingestion_worker() -> IngestionWorker:
    """Get the global ingestion worker instance."""
    return ingestion_worker